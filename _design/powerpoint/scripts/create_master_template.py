#!/usr/bin/env python3
"""
Create Leland Master PPTX Template - Clean Professional Design

Designed specifically for PowerPoint, not as an HTML imitation.
Focus: Clean typography, generous whitespace, professional polish.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

# =============================================================================
# BRAND COLORS
# =============================================================================
class Colors:
    # Primary
    LELAND_GREEN = RGBColor(0x15, 0xB0, 0x78)
    DARK_GREEN = RGBColor(0x11, 0x3A, 0x2D)
    MEDIUM_GREEN = RGBColor(0x18, 0x54, 0x40)

    # Neutrals
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE = RGBColor(0xF8, 0xF8, 0xF8)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
    BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
    TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
    TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
    BLACK = RGBColor(0x1A, 0x1A, 0x1A)

    # Accents
    RED = RGBColor(0xE8, 0x4C, 0x3D)
    ORANGE = RGBColor(0xF0, 0x93, 0x2B)

    # Tints (for backgrounds)
    GREEN_TINT = RGBColor(0xE6, 0xF7, 0xF0)
    RED_TINT = RGBColor(0xFD, 0xEE, 0xEC)


# =============================================================================
# LAYOUT CONSTANTS
# =============================================================================
SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
SLIDE_HEIGHT = Inches(7.5)

# Generous margins for clean look
MARGIN = Inches(0.75)
CONTENT_LEFT = MARGIN
CONTENT_TOP = Inches(1.5)  # Room for header
CONTENT_WIDTH = SLIDE_WIDTH - (2 * MARGIN)
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - MARGIN

# Typography
FONT_TITLE = 'Calibri Light'  # Clean, modern (widely available)
FONT_BODY = 'Calibri'
FONT_MONO = 'Consolas'


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_shape_transparency(shape, percent):
    """Set fill transparency (0-100)."""
    fill = shape.fill
    fill_elem = fill._xPr
    if fill_elem is not None:
        solid_fill = fill_elem.find(qn('a:solidFill'))
        if solid_fill is not None:
            srgb = solid_fill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = srgb.find(qn('a:alpha'))
                if alpha is None:
                    from lxml import etree
                    alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha.set('val', str(int((100 - percent) * 1000)))


def add_textbox(slide, left, top, width, height, text,
                font_size=Pt(14), font_name=FONT_BODY, font_color=Colors.TEXT_DARK,
                bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                vertical=MSO_ANCHOR.TOP, line_spacing=1.15):
    """Add a text box with clean formatting."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = vertical

    # Clear default paragraph and set text
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.line_spacing = line_spacing

    return shape


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=None, border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle with optional fill and border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

    # Adjust corner radius (default is too round)
    shape.adjustments[0] = 0.05  # Subtle rounding

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    return shape


def add_circle(slide, left, top, size, fill_color, text=None, text_color=Colors.WHITE):
    """Add a circle, optionally with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return shape


def add_left_accent_bar(slide, top, height, color=Colors.LELAND_GREEN):
    """Add a vertical accent bar on the left side of a content area."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CONTENT_LEFT, top,
        Inches(0.08), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


# =============================================================================
# SLIDE HEADER COMPONENT
# =============================================================================

def add_slide_header(slide, eyebrow=None, title=None, eyebrow_color=Colors.LELAND_GREEN):
    """Add consistent header to content slides."""
    y = Inches(0.6)

    if eyebrow:
        add_textbox(
            slide, CONTENT_LEFT, y, Inches(6), Inches(0.35),
            eyebrow.upper(),
            font_size=Pt(11), font_color=eyebrow_color,
            bold=True, font_name=FONT_BODY
        )
        y += Inches(0.35)

    if title:
        add_textbox(
            slide, CONTENT_LEFT, y, CONTENT_WIDTH, Inches(0.6),
            title,
            font_size=Pt(36), font_color=Colors.BLACK,
            bold=False, font_name=FONT_TITLE
        )


# =============================================================================
# SLIDE LAYOUTS
# =============================================================================

def create_title_slide(prs):
    """
    SLIDE 1: Title Slide
    Dark green background, centered content, clean typography.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.DARK_GREEN

    # Eyebrow - subtle, above title
    add_textbox(
        slide, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.4),
        "GRE SELF-STUDY COURSE",
        font_size=Pt(12), font_color=Colors.LELAND_GREEN,
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # Main title - large, white
    add_textbox(
        slide, Inches(0), Inches(3.0), SLIDE_WIDTH, Inches(1.0),
        "Integers & Order of Operations",
        font_size=Pt(48), font_color=Colors.WHITE,
        bold=False, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER
    )

    # Subtitle - lighter, smaller
    add_textbox(
        slide, Inches(0), Inches(4.1), SLIDE_WIDTH, Inches(0.5),
        "Master the rules. Stop the careless errors.",
        font_size=Pt(18), font_color=RGBColor(0x99, 0x99, 0x99),
        alignment=PP_ALIGN.CENTER
    )

    return slide


def create_section_divider(prs, part_num="01", title="Section Title"):
    """
    SLIDE 2: Section Divider
    Green background, minimal, bold.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Green background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.LELAND_GREEN

    # Part number - subtle
    add_textbox(
        slide, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(0.4),
        f"PART {part_num}",
        font_size=Pt(14), font_color=Colors.WHITE,
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # Section title - large, white
    add_textbox(
        slide, Inches(1), Inches(3.3), SLIDE_WIDTH - Inches(2), Inches(1.2),
        title,
        font_size=Pt(44), font_color=Colors.WHITE,
        font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER
    )

    return slide


def create_hook_slide(prs):
    """
    SLIDE 3: Hook/Statement Slide
    Big bold statement to grab attention.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.WHITE

    # Eyebrow
    add_textbox(
        slide, CONTENT_LEFT, Inches(1.5), Inches(4), Inches(0.35),
        "WHY THIS MATTERS",
        font_size=Pt(11), font_color=Colors.LELAND_GREEN, bold=True
    )

    # Main statement - large, impactful
    add_textbox(
        slide, CONTENT_LEFT, Inches(2.0), Inches(10), Inches(1.8),
        "Arithmetic errors are the #1 reason students lose points they shouldn't on the GRE.",
        font_size=Pt(36), font_color=Colors.BLACK,
        font_name=FONT_TITLE, line_spacing=1.25
    )

    # Supporting text
    add_textbox(
        slide, CONTENT_LEFT, Inches(4.2), Inches(9), Inches(1.0),
        "Not because the math is hard, but because under pressure, signs get flipped and operations get done in the wrong order.",
        font_size=Pt(18), font_color=Colors.TEXT_GRAY, line_spacing=1.4
    )

    return slide


def create_two_rules_slide(prs):
    """
    SLIDE 4: Two Rules/Concepts Side by Side
    Clean two-column layout with cards.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.OFF_WHITE

    add_slide_header(slide, "SIGNED NUMBERS", "Addition Rules")

    # Two cards side by side
    card_width = Inches(5.8)
    card_height = Inches(4.2)
    gap = Inches(0.5)
    cards_top = Inches(1.8)

    for i, (rule_name, rule_title, desc, examples) in enumerate([
        ("RULE 1", "Same Signs", "Add the absolute values, keep the sign.",
         ["6 + 3 = 9", "(−6) + (−3) = −9"]),
        ("RULE 2", "Different Signs", "Subtract smaller from larger. Keep the larger number's sign.",
         ["6 + (−3) = 3", "(−6) + 3 = −3"])
    ]):
        x = CONTENT_LEFT + i * (card_width + gap)

        # Card background
        card = add_rounded_rect(slide, x, cards_top, card_width, card_height,
                               fill_color=Colors.WHITE, border_color=Colors.BORDER_GRAY)

        # Rule label
        add_textbox(
            slide, x + Inches(0.4), cards_top + Inches(0.35), Inches(2), Inches(0.3),
            rule_name,
            font_size=Pt(11), font_color=Colors.LELAND_GREEN, bold=True
        )

        # Rule title
        add_textbox(
            slide, x + Inches(0.4), cards_top + Inches(0.7), card_width - Inches(0.8), Inches(0.5),
            rule_title,
            font_size=Pt(24), font_color=Colors.BLACK, font_name=FONT_TITLE
        )

        # Description
        add_textbox(
            slide, x + Inches(0.4), cards_top + Inches(1.25), card_width - Inches(0.8), Inches(0.6),
            desc,
            font_size=Pt(14), font_color=Colors.TEXT_GRAY, line_spacing=1.3
        )

        # Examples
        example_top = cards_top + Inches(2.0)
        for j, ex in enumerate(examples):
            ex_box = add_rounded_rect(
                slide, x + Inches(0.4), example_top + j * Inches(0.7),
                card_width - Inches(0.8), Inches(0.55),
                fill_color=Colors.LIGHT_GRAY
            )
            add_textbox(
                slide, x + Inches(0.6), example_top + j * Inches(0.7) + Inches(0.12),
                card_width - Inches(1.2), Inches(0.35),
                ex,
                font_size=Pt(18), font_name=FONT_MONO, font_color=Colors.TEXT_DARK
            )

    return slide


def create_key_concept_slide(prs):
    """
    SLIDE 5: Key Concept with Accent Bar
    Important rule or principle with visual emphasis.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.WHITE

    add_slide_header(slide, "SIGNED NUMBERS", "Subtraction = Add the Opposite")

    # Key concept box with left accent
    box_top = Inches(2.0)
    box_height = Inches(1.0)

    # Light green background
    concept_bg = add_rounded_rect(
        slide, CONTENT_LEFT, box_top, CONTENT_WIDTH, box_height,
        fill_color=Colors.GREEN_TINT
    )

    # Green left accent bar
    add_left_accent_bar(slide, box_top, box_height)

    # Concept text
    add_textbox(
        slide, CONTENT_LEFT + Inches(0.5), box_top + Inches(0.25),
        CONTENT_WIDTH - Inches(1), Inches(0.6),
        "To subtract, add the opposite of the number being subtracted.",
        font_size=Pt(20), font_color=Colors.TEXT_DARK, bold=True
    )

    # Three example boxes below
    ex_top = Inches(3.4)
    ex_width = (CONTENT_WIDTH - Inches(0.6)) / 3

    examples = [
        ("6 − 3", "= 6 + (−3)", "= 3"),
        ("6 − (−3)", "= 6 + (+3)", "= 9"),
        ("−6 − 3", "= −6 + (−3)", "= −9"),
    ]

    for i, (line1, line2, line3) in enumerate(examples):
        x = CONTENT_LEFT + i * (ex_width + Inches(0.3))

        box = add_rounded_rect(
            slide, x, ex_top, ex_width, Inches(2.0),
            fill_color=Colors.LIGHT_GRAY
        )

        # Three lines of math
        add_textbox(
            slide, x + Inches(0.3), ex_top + Inches(0.3), ex_width - Inches(0.6), Inches(0.4),
            line1, font_size=Pt(22), font_name=FONT_MONO, font_color=Colors.TEXT_DARK
        )
        add_textbox(
            slide, x + Inches(0.3), ex_top + Inches(0.8), ex_width - Inches(0.6), Inches(0.4),
            line2, font_size=Pt(16), font_name=FONT_MONO, font_color=Colors.TEXT_GRAY
        )
        add_textbox(
            slide, x + Inches(0.3), ex_top + Inches(1.3), ex_width - Inches(0.6), Inches(0.4),
            line3, font_size=Pt(22), font_name=FONT_MONO, font_color=Colors.LELAND_GREEN, bold=True
        )

    # Memory tip at bottom
    tip_top = Inches(5.8)
    tip_bg = add_rounded_rect(
        slide, CONTENT_LEFT, tip_top, CONTENT_WIDTH, Inches(0.7),
        fill_color=Colors.DARK_GREEN
    )
    add_textbox(
        slide, CONTENT_LEFT, tip_top + Inches(0.15), CONTENT_WIDTH, Inches(0.45),
        "Memory tip: \"Subtracting a negative is adding a positive.\"",
        font_size=Pt(16), font_color=Colors.WHITE, alignment=PP_ALIGN.CENTER
    )

    return slide


def create_problem_solution_slide(prs):
    """
    SLIDE 6: Problem + Step-by-Step Solution
    Clean math problem walkthrough.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.OFF_WHITE

    add_slide_header(slide, "EXAMPLE 1", "Basic Signed Arithmetic")

    # Left side: Problem
    problem_width = Inches(4.5)

    problem_box = add_rounded_rect(
        slide, CONTENT_LEFT, CONTENT_TOP, problem_width, Inches(2.5),
        fill_color=Colors.WHITE, border_color=Colors.BORDER_GRAY
    )

    add_textbox(
        slide, CONTENT_LEFT + Inches(0.4), CONTENT_TOP + Inches(0.3),
        Inches(2), Inches(0.3),
        "PROBLEM",
        font_size=Pt(10), font_color=Colors.TEXT_GRAY, bold=True
    )

    add_textbox(
        slide, CONTENT_LEFT + Inches(0.4), CONTENT_TOP + Inches(0.8),
        problem_width - Inches(0.8), Inches(0.6),
        "−5 + (−6) − 8",
        font_size=Pt(32), font_name=FONT_MONO, font_color=Colors.TEXT_DARK
    )

    # Right side: Solution steps
    sol_left = CONTENT_LEFT + problem_width + Inches(0.5)
    sol_width = CONTENT_WIDTH - problem_width - Inches(0.5)

    add_textbox(
        slide, sol_left, CONTENT_TOP, Inches(2), Inches(0.3),
        "SOLUTION",
        font_size=Pt(10), font_color=Colors.TEXT_GRAY, bold=True
    )

    steps = [
        ("Add same signs", "−5 + (−6) = −11"),
        ("Convert subtraction", "−11 − 8 = −11 + (−8)"),
        ("Add same signs", "−11 + (−8) = −19"),
    ]

    step_top = CONTENT_TOP + Inches(0.5)
    step_height = Inches(0.8)

    for i, (title, formula) in enumerate(steps):
        y = step_top + i * (step_height + Inches(0.15))

        # Step box
        step_box = add_rounded_rect(
            slide, sol_left, y, sol_width, step_height,
            fill_color=Colors.WHITE, border_color=Colors.BORDER_GRAY
        )

        # Step number circle
        add_circle(slide, sol_left + Inches(0.25), y + Inches(0.2), Inches(0.4),
                  Colors.LELAND_GREEN, str(i + 1))

        # Step title
        add_textbox(
            slide, sol_left + Inches(0.8), y + Inches(0.15),
            sol_width - Inches(1), Inches(0.3),
            title,
            font_size=Pt(13), font_color=Colors.TEXT_DARK, bold=True
        )

        # Formula
        add_textbox(
            slide, sol_left + Inches(0.8), y + Inches(0.45),
            sol_width - Inches(1), Inches(0.3),
            formula,
            font_size=Pt(15), font_name=FONT_MONO, font_color=Colors.TEXT_GRAY
        )

    # Answer bar
    answer_top = step_top + 3 * (step_height + Inches(0.15)) + Inches(0.2)
    answer_box = add_rounded_rect(
        slide, sol_left, answer_top, sol_width, Inches(0.65),
        fill_color=Colors.LELAND_GREEN
    )
    add_textbox(
        slide, sol_left + Inches(0.3), answer_top + Inches(0.15),
        sol_width - Inches(0.6), Inches(0.4),
        "Answer: −19",
        font_size=Pt(20), font_color=Colors.WHITE, bold=True
    )

    return slide


def create_warning_slide(prs):
    """
    SLIDE 7: Warning/Trap Alert
    Red-tinted alert for common mistakes.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.WHITE

    add_slide_header(slide, "CRITICAL GRE TRAP", "The Left-to-Right Rule", Colors.RED)

    # Warning box
    warn_top = Inches(2.0)
    warn_box = add_rounded_rect(
        slide, CONTENT_LEFT, warn_top, CONTENT_WIDTH, Inches(1.3),
        fill_color=Colors.RED_TINT, border_color=Colors.RED
    )

    add_textbox(
        slide, CONTENT_LEFT + Inches(0.5), warn_top + Inches(0.35),
        CONTENT_WIDTH - Inches(1), Inches(0.8),
        "Multiplication and division have EQUAL priority.\nGo left to right — don't do all multiplication first!",
        font_size=Pt(18), font_color=Colors.TEXT_DARK, line_spacing=1.4
    )

    # Example showing wrong vs right
    ex_top = Inches(3.8)
    ex_width = Inches(5.5)

    # Wrong way
    add_textbox(
        slide, CONTENT_LEFT, ex_top, ex_width, Inches(0.4),
        "24 ÷ 4 × 2",
        font_size=Pt(28), font_name=FONT_MONO, font_color=Colors.TEXT_DARK
    )

    add_textbox(
        slide, CONTENT_LEFT, ex_top + Inches(0.6), ex_width, Inches(0.35),
        "WRONG: 24 ÷ 8 = 3",
        font_size=Pt(16), font_color=Colors.RED, bold=True
    )

    add_textbox(
        slide, CONTENT_LEFT, ex_top + Inches(1.0), ex_width, Inches(0.35),
        "RIGHT: 6 × 2 = 12",
        font_size=Pt(16), font_color=Colors.LELAND_GREEN, bold=True
    )

    # Explanation
    add_textbox(
        slide, CONTENT_LEFT + Inches(6), ex_top + Inches(0.2), Inches(5), Inches(1.5),
        "Work left to right:\n\n24 ÷ 4 = 6\n6 × 2 = 12",
        font_size=Pt(18), font_name=FONT_MONO, font_color=Colors.TEXT_GRAY, line_spacing=1.5
    )

    return slide


def create_summary_table_slide(prs):
    """
    SLIDE 8: Summary Table
    Clean table with alternating rows.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.WHITE

    add_slide_header(slide, "QUICK REFERENCE", "Sign Rules Summary")

    # Table dimensions
    table_left = CONTENT_LEFT
    table_top = Inches(2.0)
    col_widths = [Inches(2.5), Inches(4.5), Inches(4.5)]
    row_height = Inches(0.7)

    headers = ["Operation", "Same Signs", "Different Signs"]
    rows = [
        ["Add", "Add absolutes, keep sign", "Subtract, keep larger's sign"],
        ["Subtract", "Convert to addition", "Convert to addition"],
        ["Multiply", "Positive", "Negative"],
        ["Divide", "Positive", "Negative"],
    ]

    # Header row
    x = table_left
    for col, width in zip(headers, col_widths):
        cell = add_rounded_rect(slide, x, table_top, width - Inches(0.05), row_height,
                               fill_color=Colors.DARK_GREEN)
        add_textbox(
            slide, x + Inches(0.25), table_top + Inches(0.2),
            width - Inches(0.5), Inches(0.4),
            col, font_size=Pt(14), font_color=Colors.WHITE, bold=True
        )
        x += width

    # Data rows
    for row_idx, row in enumerate(rows):
        y = table_top + (row_idx + 1) * row_height + Inches(0.05)
        bg_color = Colors.WHITE if row_idx % 2 == 0 else Colors.LIGHT_GRAY
        x = table_left

        for col_idx, (cell_text, width) in enumerate(zip(row, col_widths)):
            cell = add_rounded_rect(slide, x, y, width - Inches(0.05), row_height - Inches(0.05),
                                   fill_color=bg_color)

            # Color coding for Positive/Negative
            color = Colors.TEXT_DARK
            is_bold = col_idx == 0
            if cell_text == "Positive":
                color = Colors.LELAND_GREEN
                is_bold = True
            elif cell_text == "Negative":
                color = Colors.RED
                is_bold = True

            add_textbox(
                slide, x + Inches(0.25), y + Inches(0.2),
                width - Inches(0.5), Inches(0.4),
                cell_text, font_size=Pt(14), font_color=color, bold=is_bold
            )
            x += width

    return slide


def create_three_traps_slide(prs):
    """
    SLIDE 9: Three Key Points/Traps
    Three cards with important takeaways.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.OFF_WHITE

    add_slide_header(slide, "TRAPS TO AVOID", "Three GRE Pitfalls", Colors.RED)

    # Three cards
    card_width = (CONTENT_WIDTH - Inches(0.6)) / 3
    card_height = Inches(3.5)
    card_top = Inches(2.0)

    traps = [
        ("1", "M/D Trap", "Multiplication and division have EQUAL priority. Go left to right."),
        ("2", "Negative Squared", "(−3)² = 9\nbut −3² = −9\n\nParentheses matter!"),
        ("3", "Division by Zero", "x ÷ 0 is undefined,\nnot zero.\n\nThe GRE tests this."),
    ]

    for i, (num, title, desc) in enumerate(traps):
        x = CONTENT_LEFT + i * (card_width + Inches(0.3))

        # Card
        card = add_rounded_rect(
            slide, x, card_top, card_width, card_height,
            fill_color=Colors.WHITE, border_color=Colors.RED
        )

        # Number
        add_textbox(
            slide, x + Inches(0.35), card_top + Inches(0.3),
            Inches(0.5), Inches(0.5),
            num, font_size=Pt(32), font_color=Colors.RED, bold=True
        )

        # Title
        add_textbox(
            slide, x + Inches(0.35), card_top + Inches(0.9),
            card_width - Inches(0.7), Inches(0.5),
            title, font_size=Pt(20), font_color=Colors.TEXT_DARK, bold=True
        )

        # Description
        add_textbox(
            slide, x + Inches(0.35), card_top + Inches(1.5),
            card_width - Inches(0.7), Inches(1.8),
            desc, font_size=Pt(14), font_color=Colors.TEXT_GRAY, line_spacing=1.4
        )

    # Bottom bar
    bar_top = card_top + card_height + Inches(0.3)
    bar = add_rounded_rect(
        slide, CONTENT_LEFT, bar_top, CONTENT_WIDTH, Inches(0.65),
        fill_color=Colors.DARK_GREEN
    )
    add_textbox(
        slide, CONTENT_LEFT, bar_top + Inches(0.15), CONTENT_WIDTH, Inches(0.4),
        "Next up: Put these skills to work in the practice problems",
        font_size=Pt(16), font_color=Colors.WHITE, alignment=PP_ALIGN.CENTER
    )

    return slide


def create_closing_slide(prs):
    """
    SLIDE 10: Closing/CTA Slide
    Clean, professional close.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.DARK_GREEN

    # Main message
    add_textbox(
        slide, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(0.8),
        "Ready to Practice?",
        font_size=Pt(44), font_color=Colors.WHITE,
        font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER
    )

    # Subtitle
    add_textbox(
        slide, Inches(0), Inches(3.7), SLIDE_WIDTH, Inches(0.5),
        "Apply these rules in the Order of Operations Drill",
        font_size=Pt(18), font_color=RGBColor(0x99, 0x99, 0x99),
        alignment=PP_ALIGN.CENTER
    )

    # CTA button
    btn_width = Inches(3)
    btn_left = (SLIDE_WIDTH - btn_width) / 2
    btn = add_rounded_rect(
        slide, btn_left, Inches(4.5), btn_width, Inches(0.7),
        fill_color=Colors.LELAND_GREEN
    )
    add_textbox(
        slide, btn_left, Inches(4.62), btn_width, Inches(0.5),
        "Start Drill →",
        font_size=Pt(18), font_color=Colors.WHITE, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    return slide


# =============================================================================
# MAIN
# =============================================================================

def main():
    print("Creating Leland Master PPTX Template (Clean Design)")
    print("=" * 50)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Creating slides:")

    print("  1. Title Slide")
    create_title_slide(prs)

    print("  2. Section Divider")
    create_section_divider(prs, "01", "Signed Number Rules")

    print("  3. Hook Statement")
    create_hook_slide(prs)

    print("  4. Two Rules Side-by-Side")
    create_two_rules_slide(prs)

    print("  5. Key Concept")
    create_key_concept_slide(prs)

    print("  6. Problem + Solution")
    create_problem_solution_slide(prs)

    print("  7. Warning/Trap")
    create_warning_slide(prs)

    print("  8. Summary Table")
    create_summary_table_slide(prs)

    print("  9. Three Traps")
    create_three_traps_slide(prs)

    print("  10. Closing")
    create_closing_slide(prs)

    # Save
    script_dir = Path(__file__).parent
    output_path = script_dir.parent / "slide-templates" / "leland-master-template.pptx"

    prs.save(str(output_path))

    print()
    print(f"Saved to: {output_path}")
    print()
    print("Open in PowerPoint to review the design.")


if __name__ == "__main__":
    main()
