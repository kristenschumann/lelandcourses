#!/usr/bin/env python3
"""
Generate PPTX Slides from YAML Content

This script takes structured YAML content and generates a PowerPoint presentation
using the Leland brand design system. It supports inline color formatting for
math content using a token syntax:

    {{neg:-5}}  -> Red colored text for negative numbers
    {{ans:9}}   -> Green colored text for answers
    {{hl:text}} -> Green highlighted text

Usage:
    python generate_slides.py content.yaml output.pptx
    python generate_slides.py content.yaml  # Uses content filename as output

Example YAML:
    title: "Integers & Order of Operations"
    slides:
      - type: title_dark
        eyebrow: "GRE Self-Study Course"
        title: "Integers & Order"
        title_highlight: "of Operations"
        subtitle: "Master the rules. Stop the careless errors."
"""

import yaml
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# BRAND COLORS
# =============================================================================
COLORS = {
    'leland_green': RGBColor(0x15, 0xB0, 0x78),
    'darkest_green': RGBColor(0x11, 0x3A, 0x2D),
    'medium_green': RGBColor(0x18, 0x54, 0x40),
    'off_white': RGBColor(0xF8, 0xF8, 0xF8),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x00, 0x00, 0x00),
    'dark': RGBColor(0x33, 0x33, 0x33),
    'light': RGBColor(0x70, 0x70, 0x70),
    'extra_light': RGBColor(0x9B, 0x9B, 0x9B),
    'stroke': RGBColor(0xE5, 0xE5, 0xE5),
    'red': RGBColor(0xFB, 0x5A, 0x42),
    'orange': RGBColor(0xEF, 0x85, 0x09),
    'green_tint': RGBColor(0xE8, 0xF7, 0xF1),
    'red_tint': RGBColor(0xFE, 0xED, 0xEB),
}

# Slide dimensions
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)
MARGIN_LEFT = Inches(0.83)
MARGIN_RIGHT = Inches(0.83)
MARGIN_TOP = Inches(0.58)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT


# =============================================================================
# TOKEN PARSING FOR INLINE COLORS
# =============================================================================

# Regex to match tokens like {{neg:-5}} or {{ans:9}} or {{hl:text}}
TOKEN_PATTERN = re.compile(r'\{\{(neg|ans|hl):([^}]+)\}\}')

def parse_tokens(text):
    """
    Parse text with color tokens and return a list of (text, color) tuples.

    Example:
        "{{neg:-5}} + 3 = {{ans:-2}}"
        -> [('-5', 'red'), (' + 3 = ', None), ('-2', 'green')]
    """
    segments = []
    last_end = 0

    for match in TOKEN_PATTERN.finditer(text):
        # Add text before the token
        if match.start() > last_end:
            segments.append((text[last_end:match.start()], None))

        # Add the token content with its color
        token_type = match.group(1)
        token_content = match.group(2)

        if token_type == 'neg':
            segments.append((token_content, 'red'))
        elif token_type == 'ans':
            segments.append((token_content, 'leland_green'))
        elif token_type == 'hl':
            segments.append((token_content, 'leland_green'))

        last_end = match.end()

    # Add remaining text
    if last_end < len(text):
        segments.append((text[last_end:], None))

    return segments if segments else [(text, None)]


def add_formatted_text(text_frame, text, font_size=Pt(14), default_color='dark',
                       font_name='Inter', bold=False, alignment=PP_ALIGN.LEFT):
    """
    Add text with inline color formatting to a text frame.
    Parses {{neg:}}, {{ans:}}, and {{hl:}} tokens.
    """
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment

    segments = parse_tokens(text)

    for i, (segment_text, color_name) in enumerate(segments):
        if i == 0:
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        else:
            run = paragraph.add_run()

        run.text = segment_text
        run.font.size = font_size
        run.font.name = font_name
        run.font.bold = bold or (color_name is not None)  # Bold colored text

        if color_name:
            run.font.color.rgb = COLORS[color_name]
        else:
            run.font.color.rgb = COLORS[default_color]


def add_text_box_formatted(slide, left, top, width, height, text, **kwargs):
    """Add a text box with formatted text (supports tokens)."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    vertical_anchor = kwargs.pop('vertical_anchor', MSO_ANCHOR.TOP)
    tf.vertical_anchor = vertical_anchor

    add_formatted_text(tf, text, **kwargs)
    return shape


def add_rectangle(slide, left, top, width, height, fill_color=None,
                  line_color=None, line_width=Pt(1), corner_radius=None):
    """Add a rectangle shape."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[fill_color] if isinstance(fill_color, str) else fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = COLORS[line_color] if isinstance(line_color, str) else line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()

    return shape


def add_step_circle(slide, left, top, number, size=Inches(0.25)):
    """Add a numbered step circle."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['leland_green']
    circle.line.fill.background()

    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['white']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    return circle


# =============================================================================
# SLIDE GENERATORS
# =============================================================================

def create_title_dark(prs, data):
    """Create a title slide with dark green background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['darkest_green']

    # Eyebrow
    if data.get('eyebrow'):
        add_text_box_formatted(
            slide, MARGIN_LEFT, Inches(1.8), Inches(6), Inches(0.3),
            data['eyebrow'].upper(),
            font_size=Pt(13), default_color='leland_green', bold=True
        )

    # Main title
    if data.get('title'):
        add_text_box_formatted(
            slide, MARGIN_LEFT, Inches(2.1), Inches(8), Inches(0.7),
            data['title'],
            font_size=Pt(44), default_color='white', bold=True
        )

    # Highlighted part of title
    if data.get('title_highlight'):
        add_text_box_formatted(
            slide, MARGIN_LEFT, Inches(2.65), Inches(8), Inches(0.7),
            data['title_highlight'],
            font_size=Pt(44), default_color='leland_green', bold=True
        )

    # Subtitle
    if data.get('subtitle'):
        add_text_box_formatted(
            slide, MARGIN_LEFT, Inches(3.3), Inches(6), Inches(0.4),
            data['subtitle'],
            font_size=Pt(18), default_color='extra_light'
        )

    return slide


def create_section_divider(prs, data):
    """Create a section divider slide with green background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['leland_green']

    # Part number
    part = data.get('part', '01')
    add_text_box_formatted(
        slide, Inches(0), Inches(2.0), SLIDE_WIDTH, Inches(0.3),
        f"Part {part}",
        font_size=Pt(14), default_color='white', bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Section title
    add_text_box_formatted(
        slide, Inches(0), Inches(2.35), SLIDE_WIDTH, Inches(0.8),
        data.get('title', 'Section Title'),
        font_size=Pt(48), default_color='white', bold=True,
        alignment=PP_ALIGN.CENTER
    )

    return slide


def create_hook_statement(prs, data):
    """Create a hook statement slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    # Eyebrow
    add_text_box_formatted(
        slide, MARGIN_LEFT, Inches(1.5), Inches(3), Inches(0.3),
        data.get('eyebrow', 'WHY THIS MATTERS').upper(),
        font_size=Pt(13), default_color='leland_green', bold=True
    )

    # Main statement (supports inline red emphasis)
    add_text_box_formatted(
        slide, MARGIN_LEFT, Inches(1.9), Inches(7.3), Inches(1.2),
        data.get('statement', ''),
        font_size=Pt(36), default_color='black', bold=True
    )

    # Description
    if data.get('description'):
        add_text_box_formatted(
            slide, MARGIN_LEFT, Inches(3.4), Inches(7.3), Inches(0.8),
            data['description'],
            font_size=Pt(18), default_color='light'
        )

    return slide


def create_two_column_rules(prs, data):
    """Create a two-column rules slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    # Eyebrow
    add_text_box_formatted(
        slide, MARGIN_LEFT, MARGIN_TOP, Inches(3), Inches(0.3),
        data.get('eyebrow', 'CONCEPT').upper(),
        font_size=Pt(13), default_color='leland_green', bold=True
    )

    # Title
    add_text_box_formatted(
        slide, MARGIN_LEFT, Inches(0.75), Inches(6), Inches(0.5),
        data.get('title', 'Rules'),
        font_size=Pt(32), default_color='black', bold=True
    )

    # Two columns
    gap = Inches(0.25)
    col_width = (CONTENT_WIDTH - gap) / 2
    col_top = Inches(1.4)
    col_height = Inches(3.5)

    columns = data.get('columns', [{}, {}])

    for i, col_data in enumerate(columns[:2]):
        x = MARGIN_LEFT + i * (col_width + gap)

        # Column box
        add_rectangle(slide, x, col_top, col_width, col_height,
                     fill_color='off_white', corner_radius=Inches(0.12))

        # Rule label
        if col_data.get('label'):
            add_text_box_formatted(
                slide, x + Inches(0.25), col_top + Inches(0.2),
                col_width - Inches(0.5), Inches(0.25),
                col_data['label'].upper(),
                font_size=Pt(14), default_color='leland_green', bold=True
            )

        # Title
        if col_data.get('title'):
            add_text_box_formatted(
                slide, x + Inches(0.25), col_top + Inches(0.5),
                col_width - Inches(0.5), Inches(0.35),
                col_data['title'],
                font_size=Pt(20), default_color='dark', bold=True
            )

        # Description
        if col_data.get('description'):
            add_text_box_formatted(
                slide, x + Inches(0.25), col_top + Inches(0.9),
                col_width - Inches(0.5), Inches(0.5),
                col_data['description'],
                font_size=Pt(14), default_color='light'
            )

        # Examples
        examples = col_data.get('examples', [])
        for j, example in enumerate(examples[:3]):
            y = col_top + Inches(1.5) + j * Inches(0.55)

            # Example box
            add_rectangle(slide, x + Inches(0.2), y,
                         col_width - Inches(0.4), Inches(0.45),
                         fill_color='white', corner_radius=Inches(0.06))

            add_text_box_formatted(
                slide, x + Inches(0.35), y + Inches(0.08),
                col_width - Inches(0.6), Inches(0.3),
                example,
                font_size=Pt(16), default_color='dark',
                font_name='Consolas'
            )

    return slide


def create_problem_solution(prs, data):
    """Create a problem + solution slide with step-by-step walkthrough."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['off_white']

    # Eyebrow
    add_text_box_formatted(
        slide, MARGIN_LEFT, MARGIN_TOP, Inches(3), Inches(0.3),
        data.get('eyebrow', 'EXAMPLE').upper(),
        font_size=Pt(13), default_color='leland_green', bold=True
    )

    # Title
    add_text_box_formatted(
        slide, MARGIN_LEFT, Inches(0.75), Inches(6), Inches(0.5),
        data.get('title', 'Problem'),
        font_size=Pt(28), default_color='black', bold=True
    )

    # Layout
    gap = Inches(0.35)
    left_width = Inches(4)
    right_width = CONTENT_WIDTH - left_width - gap
    content_top = Inches(1.35)

    # Problem box
    add_rectangle(slide, MARGIN_LEFT, content_top, left_width, Inches(2.2),
                 fill_color='white', line_color='stroke', corner_radius=Inches(0.12))

    add_text_box_formatted(
        slide, MARGIN_LEFT + Inches(0.25), content_top + Inches(0.15),
        Inches(1.5), Inches(0.25),
        "PROBLEM",
        font_size=Pt(12), default_color='extra_light', bold=True
    )

    add_text_box_formatted(
        slide, MARGIN_LEFT + Inches(0.25), content_top + Inches(0.5),
        left_width - Inches(0.5), Inches(0.6),
        data.get('problem', ''),
        font_size=Pt(28), default_color='dark', font_name='Consolas'
    )

    # Given section if present
    if data.get('given'):
        given_top = content_top + Inches(1.1)
        add_rectangle(slide, MARGIN_LEFT + Inches(0.2), given_top,
                     left_width - Inches(0.4), Inches(0.9),
                     fill_color='off_white', corner_radius=Inches(0.08))

        add_text_box_formatted(
            slide, MARGIN_LEFT + Inches(0.35), given_top + Inches(0.1),
            Inches(1), Inches(0.2),
            "GIVEN",
            font_size=Pt(12), default_color='extra_light', bold=True
        )

        add_text_box_formatted(
            slide, MARGIN_LEFT + Inches(0.35), given_top + Inches(0.35),
            left_width - Inches(0.7), Inches(0.5),
            data['given'],
            font_size=Pt(14), default_color='dark'
        )

    # Solution section
    right_x = MARGIN_LEFT + left_width + gap

    add_text_box_formatted(
        slide, right_x, content_top, Inches(2), Inches(0.25),
        "SOLUTION",
        font_size=Pt(12), default_color='extra_light', bold=True
    )

    # Steps
    steps = data.get('steps', [])
    step_height = Inches(0.6)
    step_gap = Inches(0.1)
    steps_start = content_top + Inches(0.35)

    for i, step in enumerate(steps[:4]):
        y = steps_start + i * (step_height + step_gap)

        add_rectangle(slide, right_x, y, right_width, step_height,
                     fill_color='white', line_color='stroke', corner_radius=Inches(0.08))

        add_step_circle(slide, right_x + Inches(0.12), y + Inches(0.15), i + 1)

        step_title = step.get('title', f'Step {i+1}')
        step_formula = step.get('formula', '')

        add_text_box_formatted(
            slide, right_x + Inches(0.5), y + Inches(0.08),
            right_width - Inches(0.6), Inches(0.2),
            step_title,
            font_size=Pt(13), default_color='dark', bold=True
        )

        add_text_box_formatted(
            slide, right_x + Inches(0.5), y + Inches(0.32),
            right_width - Inches(0.6), Inches(0.25),
            step_formula,
            font_size=Pt(14), default_color='light', font_name='Consolas'
        )

    # Answer bar
    answer_y = steps_start + len(steps) * (step_height + step_gap) + Inches(0.1)
    add_rectangle(slide, right_x, answer_y, right_width, Inches(0.5),
                 fill_color='leland_green', corner_radius=Inches(0.08))

    answer = data.get('answer', '')
    add_text_box_formatted(
        slide, right_x + Inches(0.15), answer_y + Inches(0.1),
        right_width - Inches(0.3), Inches(0.35),
        f"Answer: {answer}",
        font_size=Pt(18), default_color='white', bold=True
    )

    return slide


def create_summary_table(prs, data):
    """Create a summary table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    # Eyebrow
    add_text_box_formatted(
        slide, MARGIN_LEFT, MARGIN_TOP, Inches(3), Inches(0.3),
        data.get('eyebrow', 'QUICK REFERENCE').upper(),
        font_size=Pt(13), default_color='leland_green', bold=True
    )

    # Title
    add_text_box_formatted(
        slide, MARGIN_LEFT, Inches(0.75), Inches(6), Inches(0.5),
        data.get('title', 'Summary'),
        font_size=Pt(32), default_color='black', bold=True
    )

    # Table
    table_top = Inches(1.35)
    row_height = Inches(0.55)
    headers = data.get('headers', ['Column 1', 'Column 2', 'Column 3'])
    rows = data.get('rows', [])

    col_count = len(headers)
    col_width = CONTENT_WIDTH / col_count

    # Header row
    for i, header in enumerate(headers):
        x = MARGIN_LEFT + i * col_width
        add_rectangle(slide, x, table_top, col_width - Inches(0.02), row_height,
                     fill_color='darkest_green')

        add_text_box_formatted(
            slide, x + Inches(0.2), table_top + Inches(0.15),
            col_width - Inches(0.4), Inches(0.3),
            header,
            font_size=Pt(14), default_color='white', bold=True
        )

    # Data rows
    for row_idx, row_data in enumerate(rows[:6]):
        y = table_top + row_height + Inches(0.02) + row_idx * (row_height + Inches(0.02))
        bg_color = 'white' if row_idx % 2 == 0 else 'off_white'

        for col_idx, cell in enumerate(row_data[:col_count]):
            x = MARGIN_LEFT + col_idx * col_width
            add_rectangle(slide, x, y, col_width - Inches(0.02), row_height,
                         fill_color=bg_color)

            # Determine text color
            text_color = 'dark'
            is_bold = col_idx == 0

            # Check for special color keywords
            cell_lower = cell.lower() if isinstance(cell, str) else ''
            if 'positive' in cell_lower:
                text_color = 'leland_green'
                is_bold = True
            elif 'negative' in cell_lower:
                text_color = 'red'
                is_bold = True

            add_text_box_formatted(
                slide, x + Inches(0.2), y + Inches(0.15),
                col_width - Inches(0.4), Inches(0.3),
                cell,
                font_size=Pt(14), default_color=text_color, bold=is_bold
            )

    return slide


def create_three_card_summary(prs, data):
    """Create a three-card summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']

    # Eyebrow
    eyebrow_color = 'red' if data.get('type_alert') else 'leland_green'
    add_text_box_formatted(
        slide, MARGIN_LEFT, MARGIN_TOP, Inches(3), Inches(0.3),
        data.get('eyebrow', 'SUMMARY').upper(),
        font_size=Pt(13), default_color=eyebrow_color, bold=True
    )

    # Title
    add_text_box_formatted(
        slide, MARGIN_LEFT, Inches(0.75), Inches(6), Inches(0.5),
        data.get('title', 'Key Points'),
        font_size=Pt(32), default_color='black', bold=True
    )

    # Cards
    cards = data.get('cards', [])
    gap = Inches(0.2)
    card_width = (CONTENT_WIDTH - 2 * gap) / 3
    card_top = Inches(1.35)
    card_height = Inches(2.8)

    card_bg = 'red_tint' if data.get('type_alert') else 'green_tint'
    card_border = 'red' if data.get('type_alert') else 'leland_green'
    num_color = 'red' if data.get('type_alert') else 'leland_green'

    for i, card in enumerate(cards[:3]):
        x = MARGIN_LEFT + i * (card_width + gap)

        add_rectangle(slide, x, card_top, card_width, card_height,
                     fill_color=card_bg, line_color=card_border,
                     line_width=Pt(1), corner_radius=Inches(0.12))

        # Number
        add_text_box_formatted(
            slide, x + Inches(0.25), card_top + Inches(0.2),
            Inches(0.5), Inches(0.5),
            str(i + 1),
            font_size=Pt(28), default_color=num_color, bold=True
        )

        # Title
        add_text_box_formatted(
            slide, x + Inches(0.25), card_top + Inches(0.75),
            card_width - Inches(0.5), Inches(0.4),
            card.get('title', ''),
            font_size=Pt(18), default_color='dark', bold=True
        )

        # Description
        add_text_box_formatted(
            slide, x + Inches(0.25), card_top + Inches(1.25),
            card_width - Inches(0.5), Inches(1.2),
            card.get('description', ''),
            font_size=Pt(14), default_color='light'
        )

    # Bottom bar if present
    if data.get('bottom_text'):
        bar_y = card_top + card_height + Inches(0.2)
        add_rectangle(slide, MARGIN_LEFT, bar_y, CONTENT_WIDTH, Inches(0.5),
                     fill_color='darkest_green', corner_radius=Inches(0.08))

        add_text_box_formatted(
            slide, MARGIN_LEFT + Inches(0.2), bar_y + Inches(0.1),
            CONTENT_WIDTH - Inches(0.4), Inches(0.35),
            data['bottom_text'],
            font_size=Pt(16), default_color='white',
            alignment=PP_ALIGN.CENTER
        )

    return slide


# =============================================================================
# SLIDE TYPE MAPPING
# =============================================================================

SLIDE_CREATORS = {
    'title_dark': create_title_dark,
    'section_divider': create_section_divider,
    'hook_statement': create_hook_statement,
    'two_column_rules': create_two_column_rules,
    'problem_solution': create_problem_solution,
    'summary_table': create_summary_table,
    'three_card_summary': create_three_card_summary,
}


def generate_presentation(content_path, output_path=None):
    """Generate a PPTX from YAML content file."""
    content_path = Path(content_path)

    if output_path is None:
        output_path = content_path.with_suffix('.pptx')
    else:
        output_path = Path(output_path)

    # Load content
    with open(content_path, 'r', encoding='utf-8') as f:
        content = yaml.safe_load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Process slides
    slides = content.get('slides', [])

    for i, slide_data in enumerate(slides):
        slide_type = slide_data.get('type', 'title_dark')

        if slide_type in SLIDE_CREATORS:
            print(f"  Creating slide {i+1}: {slide_type}")
            SLIDE_CREATORS[slide_type](prs, slide_data)
        else:
            print(f"  Warning: Unknown slide type '{slide_type}' at position {i+1}")

    # Save
    prs.save(str(output_path))
    print(f"\nSaved to: {output_path}")

    return output_path


def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_slides.py content.yaml [output.pptx]")
        print("\nExample YAML structure:")
        print("""
title: "Presentation Title"
slides:
  - type: title_dark
    eyebrow: "Course Name"
    title: "Main Title"
    title_highlight: "Highlighted Part"
    subtitle: "Subtitle text"

  - type: section_divider
    part: "01"
    title: "Section Name"

  - type: problem_solution
    eyebrow: "Example 1"
    title: "Problem Title"
    problem: "{{neg:-5}} + ({{neg:-6}}) - 8"
    steps:
      - title: "Add same signs"
        formula: "{{neg:-5}} + ({{neg:-6}}) = {{neg:-11}}"
      - title: "Convert subtraction"
        formula: "{{neg:-11}} - 8 = {{neg:-11}} + ({{neg:-8}})"
    answer: "{{neg:-19}}"
""")
        sys.exit(1)

    content_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    print(f"Generating slides from: {content_path}")
    generate_presentation(content_path, output_path)


if __name__ == "__main__":
    main()
